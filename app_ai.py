import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import io
import re
import textwrap
import docx
from io import BytesIO
from sentence_transformers import SentenceTransformer, util

# Streamlit 세팅
st.set_page_config(page_title="Paydo AI PPT", layout="centered")
# st.title("🎬 AI PPT 생성기 (KoSimCSE)") # 이 부분 제거

# CSS 스타일 정의
custom_css = """
<style>
    /* 기본 폰트 설정 (Google Noto Sans KR 폰트 임포트) */
    @import url('https://fonts.googleapis.com/css2?family=Noto+Sans+KR:wght@400;700&display=swap');
    /* Font Awesome 아이콘 라이브러리 임포트 */
    @import url('https://cdnjs.cloudflare.com/ajax/libs/font-awesome/5.15.4/css/all.min.css');

    /* Streamlit 앱의 전체적인 배경 및 폰트 설정 */
    html, body, [data-testid="stAppViewContainer"] {
        font-family: 'Noto Sans KR', sans-serif;
        margin: 0;
        padding: 0;
        background-color: #f0f2f5; /* 전체 앱 배경색 */
        color: #333; /* 기본 텍스트 색상 */
    }

    /* Streamlit 메인 컨테이너 폭 조절 및 그림자, 모서리 둥글게 */
    /* st.set_page_config(layout="centered") 사용 시, margin, left, transition 등은 Streamlit이 자체적으로 관리하므로 제거 */
    [data-testid="stAppViewContainer"] {
        max-width: 1000px; /* 컨테이너 최대 너비 유지 */
        margin: auto; /* 페이지 중앙 정렬 (Streamlit이 관리하지만 명시적으로 유지) */
        box-shadow: 0 4px 12px rgba(0, 0, 0, 0.1); /* 그림자 효과 */
        border-radius: 8px; /* 모서리 둥글게 */
        overflow-y: auto; /* 내용이 넘칠 때 세로 스크롤 허용 (이전 요청 반영) */
        overflow-x: hidden; /* 가로 스크롤 방지 */
        background-color: #fff; /* 메인 컨테이너 배경색을 흰색으로 설정 */
        
        /* 하단 고정 바 때문에 메인 컨테이너 하단에 패딩 추가 */
        padding-bottom: 90px; /* 하단 고정 바의 높이(padding 15+15+버튼 높이 고려)에 맞춰 조절 */
    }

    /* 상단 디자인 BAR 스타일 (기존 stHeader 오버라이드 대신 직접 마크다운에 적용) */
    .top-design-bar {
        background-color: #2c3e50; /* 어두운 파란색/회색 */
        color: #fff;
        padding: 15px 20px;
        text-align: center;
        border-top-left-radius: 8px;
        border-top-right-radius: 8px;
        box-shadow: 0 2px 5px rgba(0, 0, 0, 0.1);
        /* Streamlit 기본 block-container의 패딩을 덮기 위해 음수 마진 사용 */
        margin-left: -1rem; 
        margin-right: -1rem;
        width: calc(100% + 2rem); /* 부모 너비에 맞춰 확장 */
        position: sticky; /* 스크롤 시 상단에 고정 */
        top: 0; /* 상단에 붙임 */
        z-index: 999; /* 다른 요소 위에 표시되도록 */
    }
    .top-design-bar h1 {
        color: #fff; /* 제목 텍스트 색상 흰색 */
        margin: 0;
        font-size: 1.8em; /* 제목 글자 크기 조정 (0.1em이 너무 작았음, 적절한 크기로 수정) */
        font-weight: 700;
        text-align: center; /* 가운데 정렬 */
        display: flex; /* flexbox 사용 (이모지와 텍스트 정렬) */
        align-items: center;
        justify-content: center;
        gap: 10px; /* 이모지와 텍스트 사이 간격 */
    }

    /* 고정된 하단 바 스타일 (새로 추가) */
    .bottom-fixed-bar { 
        background-color: #A2D9CE; /* 옅은 녹색으로 변경 */
        padding: 15px 20px;
        text-align: center;
        box-shadow: 0 -2px 5px rgba(0, 0, 0, 0.1);
        position: fixed; /* 뷰포트 하단에 고정 */
        bottom: 0; /* 하단에 붙임 */
        left: 50%; /* 왼쪽 50% 이동 */
        transform: translateX(-50%); /* 자신의 너비의 절반만큼 왼쪽으로 이동하여 중앙 정렬 */
        width: 100%; /* 너비 100% */
        max-width: 1000px; /* 메인 컨테이너와 동일한 최대 너비 적용 */
        z-index: 1000; /* 다른 요소 위에 표시되도록 가장 높은 z-index 부여 */
        box-sizing: border-box; /* padding이 width에 포함되도록 */
        border-bottom-left-radius: 8px; /* 메인 컨테이너와 일치하도록 */
        border-bottom-right-radius: 8px; /* 메인 컨테이너와 일치하도록 */
    }

    /* 하단 고정 바 안에 있는 실제 버튼 (button 태그) 스타일 */
    .bottom-fixed-bar .stButton > button { 
        background-color: #007BFF; /* 눈에 띄는 파란색으로 변경 */
        color: white;
        border: none;
        padding: 12px 25px; /* 패딩 증가로 버튼 크기 키우기 */
        border-radius: 8px; /* 더 둥글게 */
        cursor: pointer;
        font-size: 1.3em; /* 폰트 크기 키우기 */
        font-weight: 700;
        width: auto; /* 버튼 콘텐츠 크기에 맞게 너비 조절 */
        max-width: 400px; /* 최대 너비 제한 (너무 길어지는 것을 방지) */
        display: flex; /* flexbox 사용 */
        align-items: center;
        justify-content: center;
        gap: 10px;
        transition: background-color 0.3s ease;
    }
    .bottom-fixed-bar .stButton > button:hover {
        background-color: #0056b3; /* 호버 시 더 어두운 파란색 */
    }

    /* 대본 입력 방식 선택 섹션 */
    .input-method-selection-box {
        background-color: #e0f2f7; /* 연한 파란색 배경 */
        padding: 10px 15px;
        border-radius: 8px;
        margin-bottom: 20px;
        text-align: center;
        display: flex;
        justify-content: center;
        align-items: center;
        gap: 8px;
        font-weight: 700;
        color: #2c3e50;
        font-size: 1.1em;
    }
    .input-method-selection-box .icon {
        font-size: 1.2em;
    }

    /* Streamlit 탭 위젯 커스터마이징 */
    .stTabs [data-baseweb="tab-list"] {
        gap: 0px;
        border-bottom: 1px solid #ddd;
        margin-bottom: 20px;
    }
    .stTabs [data-baseweb="tab"] {
        background-color: lightcyan; /* 파스텔 톤 옅은 색 */
        border-radius: 4px 4px 0px 0px;
        padding: 10px 15px;
        font-weight: 500;
        color: #555;
    }
    .stTabs [aria-selected="true"] { 
        border-bottom: 2px solid #3498db !important; 
        color: #3498db !important; 
        font-weight: 700;
        background-color: lightblue; /* 선택된 탭 더 진한 파스텔 톤 */
    }
    .stTabs [data-baseweb="tab"]:hover {
        background-color: #f5f5f5;
    }

    /* Streamlit 파일 업로더 커스터마이징 */
    [data-testid="stFileUploaderDropzone"] {
        border: 2px dashed #a0d8f0;
        border-radius: 8px;
        background-color: #f7fcfe;
        padding: 30px 20px;
        height: 250px; /* 높이 고정 - 유지 */
        display: flex;
        flex-direction: column;
        justify-content: center;
        align-items: center;
        position: relative; /* 자식 요소 절대 위치 지정을 위해 */
    }
    /* 파일 업로더의 기본 텍스트와 아이콘 숨기기 */
    [data-testid="stFileUploaderDropzoneInstructions"] > div > span {
        display: none; 
    }
    [data-testid="stFileUploaderDropzoneInstructions"] > div > small {
        display: none; 
    }
    [data-testid="stFileUploaderDropzoneInstructions"] {
        display: none; 
    }
    
    /* Browse files 버튼 스타일 조정 */
    [data-testid="stFileUploaderBrowseButton"] > button {
        background-color: #3498db;
        color: white;
        border: none;
        padding: 10px 20px;
        border-radius: 5px;
        cursor: pointer;
        font-size: 0.9em;
        font-weight: 600;
        transition: background-color 0.3s ease;
        position: absolute; /* 드롭존 내에서 절대 위치 지정 */
        bottom: 20px;
        right: 20px;
    }
    [data-testid="stFileUploaderBrowseButton"] > button:hover {
        background-color: #2980b9;
    }

    /* 드래그 앤 드롭 영역 내 커스텀 텍스트 및 아이콘 위치 및 크기 조정 */
    .custom-file-uploader-content {
        text-align: center; 
        position: absolute; 
        bottom: 70px; /* 'Browse files' 버튼 위로 위치 조정 */
        left: 50%; 
        transform: translateX(-50%); /* 가로 중앙 정렬 */
        pointer-events: none; 
        z-index: 1;
    }
    .custom-file-uploader-content .fas {
        font-size: 2.5em; /* 아이콘 크기 줄임 */
        color: #3498db; 
        position: relative; /* 아이콘 자체를 상대적으로 이동 */
        top: -15px; /* 아이콘을 위로 15px 이동 */
        margin-bottom: 0; /* 기존 마진 제거 */
    }
    .custom-file-uploader-content p:nth-of-type(1) { /* "Drag and drop file here" */
        margin:0; 
        font-size: 1.0em; /* 폰트 크기 줄임 */
        color: #666;
    }
    .custom-file-uploader-content p:nth-of-type(2) { /* "Limit 200MB per file • DOCX" */
        margin:0; 
        font-size: 0.8em; /* 폰트 크기 줄임 */
        color: #888; 
        margin-top: 5px;
    }


    /* 문제 해결 Expander (st.expander) 스타일 */
    .stExpander {
        border: 1px solid #eee;
        border-radius: 8px;
        background-color: #f9f9f9;
        margin-top: 20px;
    }
    .stExpander > div > div > details > summary {
        color: #666;
        font-size: 0.9em;
        padding: 10px 15px;
        outline: none;
    }
    .stExpander > div > div > details > summary::marker {
        content: '';
    }
    .stExpander > div > div > details > summary::before {
        content: '▼';
        font-size: 0.8em;
        margin-right: 5px;
        transition: transform 0.2s;
    }
    .stExpander > div > div > details[open] > summary::before {
        transform: rotate(180deg);
    }
    .stExpander > div > div > details > div {
        padding: 5px 15px 10px;
        border-top: 1px dashed #eee;
        font-size: 0.85em;
        color: #777;
    }

    /* 사이드바 스타일 */
    [data-testid="stSidebar"] {
        background-color: #e7eff6; /* 사이드바 배경색 */
        border-right: 1px solid #ddd;
        box-shadow: 2px 0 5px rgba(0,0,0,0.05);
        height: 100%;
        z-index: 1000; /* 다른 요소 위에 표시 */
        padding-top: 100px; /* 사이드바 상단 여백을 늘려 숨김 버튼 위치 조절 */
    }
    /* 사이드바 내부 요소에 대한 스타일 */
    [data-testid="stSidebar"] h1, [data-testid="stSidebar"] h2 {
        margin-bottom: 0.5em; /* 제목 및 헤더 아래 여백 */
    }
    [data-testid="stSidebar"] label {
        font-weight: 600; /* 라벨 볼드 처리 */
        margin-bottom: 0.2em; /* 라벨 아래 여백 */
    }

    /* 반응형 디자인 */
    @media (max-width: 768px) {
        [data-testid="stSidebar"] {
            position: relative; /* 모바일에서는 고정 해제 */
            height: auto;
            width: 100%;
            padding-top: 0; /* 모바일에서는 패딩 초기화 */
            border-right: none;
            box-shadow: none;
        }
        [data-testid="stAppViewContainer"] {
            max-width: 100%; /* 모바일에서는 최대 너비 제거 */
            border-radius: 0; /* 모바일에서는 모서리 둥글게 처리 제거 */
            box-shadow: none; /* 모바일에서는 그림자 제거 */
            margin-left: 0 !important; /* 모바일에서는 마진 제거 */
            padding-bottom: 90px; /* 모바일에서도 하단 고정 바 패딩 유지 */
        }
        .top-design-bar, .bottom-fixed-bar {
            border-radius: 0;
        }
        .bottom-fixed-bar .stButton > button {
             width: auto; /* 모바일에서도 너비 자동 조절 */
             max-width: none; /* 모바일에서는 최대 너비 제한 해제 */
        }
    }
</style>
"""

# Streamlit 앱에 사용자 정의 CSS 주입
st.markdown(custom_css, unsafe_allow_html=True)

# 모델 로딩 (한 번만)
@st.cache_resource
def load_model():
    return SentenceTransformer("jhgan/ko-sbert-nli")

model = load_model()

# Word 파일 텍스트 추출
def extract_text_from_word(uploaded_file):
    try:
        file_bytes = BytesIO(uploaded_file.read())
        doc = docx.Document(file_bytes)
        return [p.text for p in doc.paragraphs if p.text.strip()]
    except Exception as e:
        st.error(f"Word 파일 처리 오류: {e}")
        return None

# 텍스트 줄 수 계산
def calculate_text_lines(text, max_chars_per_line):
    lines = 0
    paragraphs = text.split('\n')
    for paragraph in paragraphs:
        if not paragraph:
            lines += 1
        else:
            lines += len(textwrap.wrap(paragraph, width=max_chars_per_line, break_long_words=True))
    return lines

# 문장 분할
def smart_sentence_split(text):
    paragraphs = text.split('\n')
    sentences = []
    for paragraph in paragraphs:
        # 서술어 단독 분리 방지를 위해 문장 끝 마침표 기준이 아닌, 약간 넓게 split
        # 한글 문장 분리 시 '.(마침표)' 뒤에 한글이 오는 경우 오류 발생 방지
        temp_sentences = re.split(r'(?<=[^\d][.!?])\s+(?=[\"\'\uAC00-\D7A3])', paragraph)
        sentences.extend([s.strip() for s in temp_sentences if s.strip()])
    return sentences

# 슬라이드 분할 with 유사도 + 짧은 문장 병합 개선
def split_text_into_slides_with_similarity(text_paragraphs, max_lines_per_slide, max_chars_per_line_ppt, model, similarity_threshold=0.85):
    slides, split_flags, slide_number = [], [], 1
    current_text, current_lines, needs_check = "", 0, False

    for paragraph in text_paragraphs:
        sentences = smart_sentence_split(paragraph)
        if not sentences:
            continue

        embeddings = model.encode(sentences)

        i = 0
        while i < len(sentences):
            sentence = sentences[i]
            sentence_lines = calculate_text_lines(sentence, max_chars_per_line_ppt)

            # 다음 문장과 병합을 시도 (너무 짧은 문장 방지)
            if sentence_lines <= 2 and i + 1 < len(sentences):
                next_sentence = sentences[i + 1]
                merged = sentence + " " + next_sentence
                merged_lines = calculate_text_lines(merged, max_chars_per_line_ppt)
                if merged_lines <= max_lines_per_slide:
                    sentence = merged
                    sentence_lines = merged_lines
                    i += 1  # 추가로 하나 더 소비

            if current_lines + sentence_lines <= max_lines_per_slide:
                current_text += sentence + "\n"
                current_lines += sentence_lines
            else:
                slides.append(current_text.strip())
                split_flags.append(needs_check)
                slide_number += 1
                current_text = sentence + "\n"
                current_lines = sentence_lines
                needs_check = False
            i += 1

    if current_text:
        slides.append(current_text.strip())
        split_flags.append(needs_check)

    return slides, split_flags

def create_ppt(slide_texts, split_flags, max_chars_per_line_in_ppt=18, font_size=54):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    total_slides = len(slide_texts)

    for i, text in enumerate(slide_texts):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_text_to_slide(slide, text, font_size, PP_ALIGN.CENTER, max_chars_per_line_in_ppt)
        if split_flags[i]:
            add_check_needed_shape(slide)
        if i == total_slides - 1:
            add_end_mark(slide)
    return prs

def add_text_to_slide(slide, text, font_size, alignment, max_chars_per_line):
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(6.2))
    text_frame = textbox.text_frame
    text_frame.clear()
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    text_frame.word_wrap = True

    wrapped_lines = textwrap.wrap(text, width=max_chars_per_line, break_long_words=True)
    for line in wrapped_lines:
        p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = 'Noto Color Emoji' # Noto Sans KR 폰트 추가 설치 필요 시 고려
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.alignment = alignment
        p.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    text_frame.auto_size = None

def add_check_needed_shape(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(2.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 0)
    shape.line.color.rgb = RGBColor(0, 0, 0)
    p = shape.text_frame.paragraphs[0]
    p.text = "확인 필요!"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)
    shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p.alignment = PP_ALIGN.CENTER

def add_end_mark(slide):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(10), Inches(6), Inches(2), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 0, 0)
    shape.line.color.rgb = RGBColor(0, 0, 0)
    p = shape.text_frame.paragraphs[0]
    p.text = "끝"
    p.font.size = Pt(36)
    p.font.color.rgb = RGBColor(255, 255, 255)
    shape.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p.alignment = PP_ALIGN.CENTER

# --- Streamlit 앱 UI 구성 시작 ---

# 좌측 사이드바 (st.sidebar)
with st.sidebar:
    st.markdown("---") # 구분선 유지
    st.header("⚙️ PPT 생성 옵션") # 'PPT 설정' -> '⚙️ PPT 생성 옵션' (이모지 추가)
    # 안내 문구 수정
    st.markdown("<p style='font-size:0.9em; color:#555;'>생성될 PPT의 세부 옵션을 설정할 수 있습니다.</p>", unsafe_allow_html=True)
    
    # 슬라이드 수 설정 (이모지 추가)
    max_lines = st.slider("📏 슬라이드당 최대 줄 수", 1, 10, 4, key='sidebar_max_lines')
    # 한 줄당 최대 글자 수 (이모지 추가)
    max_chars = st.slider("🔠 한 줄당 최대 글자 수", 10, 100, 18, key='sidebar_max_chars')
    # 폰트 크기 (이모지 추가)
    font_size = st.slider("✍️ 폰트 크기", 10, 60, 54, key='sidebar_font_size')
    # 문맥 유사도 기준 (이모지 추가)
    sim_threshold = st.slider("💡 문맥 유사도 기준", 0.0, 1.0, 0.85, step=0.05, key='sidebar_sim_threshold')

    st.markdown("---")


# 상단 디자인 BAR (st.title 대신 직접 마크다운 사용)
with st.container():
    st.markdown('<div class="top-design-bar">', unsafe_allow_html=True)
    st.markdown("<h1>🎬 촬영 대본 PPT 자동 생성 AI (KoSimCSE)</h1>", unsafe_allow_html=True)
    st.markdown('</div>', unsafe_allow_html=True)

# 대본 입력 방식 선택 섹션 (더 작게, 이모지 반영)
st.markdown('<div class="input-method-selection-box"><span class="icon">📁</span> 대본 입력 방식 선택</div>', unsafe_allow_html=True)

# 탭 메뉴 구성 (st.tabs 위젯 사용)
tab1, tab2 = st.tabs(["📘 Word 파일 업로드", "📝 텍스트 직접 입력"])

uploaded_file_tab1 = None 
text_input_tab2 = ""

with tab1:
    st.write("Word 파일 (.docx)을 업로드해주세요.")

    # 파일 업로더 위젯
    uploaded_file_tab1 = st.file_uploader(
        "Upload your DOCX file here", # 이 텍스트는 내부적으로 사용되지만, CSS로 숨김.
        type=["docx"], # 허용되는 파일 형식
        accept_multiple_files=False, # 단일 파일만 허용
        label_visibility="collapsed" # 기본 라벨 숨기기
    )
    
    # 드래그 앤 드롭 영역 내 커스텀 텍스트 및 아이콘 (CSS로 위치 조정)
    st.markdown("""
        <div class="custom-file-uploader-content">
            <i class="fas fa-cloud-upload-alt"></i>
            <p>Drag and drop file here</p>
            <p>Limit 200MB per file • DOCX</p>
        </div>
    """, unsafe_allow_html=True)

    if uploaded_file_tab1 is not None:
        st.success(f"파일 '{uploaded_file_tab1.name}'이(가) 업로드되었습니다.")

    # 문제 해결 드롭다운 (st.expander 위젯 사용)
    with st.expander("🙁 Word 파일 업로드 시 문제가 발생하나요?"):
        st.write("문제가 발생할 경우 다음 사항을 확인해주세요:")
        st.markdown("- 파일 형식이 `.docx`인지 확인해주세요.")
        st.markdown("- 파일 크기가 200MB를 초과하지 않는지 확인해주세요.")
        st.markdown("- 네트워크 연결이 안정적인지 확인해주세요.")
        st.markdown("- **서버의 오류로 파일명이 한글인 경우 오류가 발생할 수 있습니다. 파일명을 영문으로 수정하여 다시 시도해주세요.**") # [NEW] 한글 파일명 오류 안내 추가
        st.markdown("- 다른 이름으로 저장 후 다시 시도해보세요.")

with tab2:
    text_input_tab2 = st.text_area(
        "대본을 직접 입력하세요.",
        height=200, 
        placeholder="여기에 대본을 입력해주세요...",
        help="여기에 입력된 텍스트로 PPT 대본이 생성됩니다."
    )

# 고정된 하단 바
st.markdown('<div class="bottom-fixed-bar">', unsafe_allow_html=True) 

# st.columns를 사용하여 버튼을 가운데 정렬
col1, col2, col3 = st.columns([1, 2, 1]) # 1:2:1 비율로 컬럼 생성 (가운데 컬럼이 넓음)
with col2: # 가운데 컬럼에 버튼 배치
    if st.button("🚀 PPT 자동 생성 시작", use_container_width=True): # use_container_width=True를 사용하여 컬럼 너비에 맞춤
        paragraphs = []
        target_file = None
        target_text_input = ""

        if uploaded_file_tab1 is not None:
            paragraphs = extract_text_from_word(uploaded_file_tab1)
        elif text_input_tab2.strip():
            paragraphs = [p.strip() for p in text_input_tab2.split("\n\n") if p.strip()]
        else:
            st.warning("PPT 생성을 위해 Word 파일을 업로드하거나 대본을 직접 입력해주세요.")
            st.stop()

        if not paragraphs:
            st.error("유효한 텍스트가 없습니다.")
            st.stop()

        with st.spinner("PPT 생성 중..."):
            slides, flags = split_text_into_slides_with_similarity(
                paragraphs, max_lines, max_chars, model, similarity_threshold=sim_threshold
            )
            ppt = create_ppt(slides, flags, max_chars, font_size)

            if ppt:
                ppt_io = io.BytesIO()
                ppt.save(ppt_io)
                ppt_io.seek(0)
                st.download_button(
                    label="📥 PPT 다운로드",
                    data=ppt_io,
                    file_name="paydo_script_ai.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
                st.success(f"총 {len(slides)}개의 슬라이드가 생성되었습니다.")
                if any(flags):
                    flagged = [i+1 for i, f in enumerate(flags) if f]
                    st.warning(f"⚠️ 확인이 필요한 슬라이드: {flagged}")
st.markdown('</div>', unsafe_allow_html=True)

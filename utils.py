import kss
from sentence_transformers import util
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
import io
import docx
import re
import textwrap
import logging

def read_script_file(uploaded_file):
    """파일을 읽어 텍스트를 추출하는 함수"""
    text = ""
    if uploaded_file.type == "text/plain":
        text = io.TextIOWrapper(uploaded_file, encoding='utf-8').read()
    elif uploaded_file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        doc = docx.Document(uploaded_file)
        for p in doc.paragraphs:
            text += p.text + "\n"
    return text

def smart_sentence_split(text):
    """KSS를 사용하여 문장을 분리하는 함수"""
    return kss.split_sentences(text)

def calculate_similarity(model, sentences):
    """문장 간 유사도를 계산하는 함수"""
    embeddings = model.encode(sentences)
    similarity_matrix = util.cos_sim(embeddings, embeddings)
    return similarity_matrix

def merge_short_sentences(sentences, max_length=16):
    """짧은 문장을 다음 문장과 병합하는 함수"""
    merged_sentences = []
    temp_sentence = ""
    for sentence in sentences:
        if len(temp_sentence + sentence) < max_length:
            temp_sentence += sentence + " "
        else:
            merged_sentences.append(temp_sentence.strip() + sentence.strip())
            temp_sentence = ""
    if temp_sentence:
        merged_sentences.append(temp_sentence.strip())
    return merged_sentences

def split_into_slides(model, text, sentences_per_slide=3, similarity_threshold=0.7):
    """텍스트를 슬라이드로 분할하는 함수"""

    # 문장 분리 및 병합
    sentences = smart_sentence_split(text)
    sentences = merge_short_sentences(sentences)

    slides = []
    current_slide = []
    for i in range(0, len(sentences), sentences_per_slide):
        slide_sentences = sentences[i:i + sentences_per_slide]
        
        # 유사도 검사 추가 (첫 슬라이드가 아니면)
        if slides and len(current_slide) > 0:
            prev_slide_sentences = slides[-1].split("\n")
            combined_sentences = prev_slide_sentences + slide_sentences
            similarity_matrix = calculate_similarity(model, combined_sentences)
            
            # 이전 슬라이드 마지막 문장과 현재 슬라이드 첫 문장의 유사도 비교
            if similarity_matrix[-1][len(prev_slide_sentences)] < similarity_threshold:
                slides.append("\n".join(current_slide))
                current_slide = []

        current_slide.append("\n".join(slide_sentences))

    slides.append("\n".join(current_slide))
    return slides

def process_script(text, model):
    """전체 대본 처리 함수"""
    slides = split_into_slides(model, text)
    slides_data = []
    for slide_content in slides:
        slides_data.append({
            "text": slide_content,
            "flags": []  # 추가: 슬라이드별 플래그 초기화
        })
    return slides_data

def create_ppt(slides_data):
    """PPT를 생성하는 함수"""
    prs = Presentation()
    for i, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 레이아웃 사용

        # 텍스트 박스 추가 및 설정 (여백, 폰트, 크기 등)
        left = top = Inches(1)
        width = prs.slide_width - Inches(2)
        height = prs.slide_height - Inches(2)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.clear()
        text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        text_frame.word_wrap = True

        # 텍스트 추가 및 가운데 정렬
        p = text_frame.add_paragraph()
        p.text = slide_data["text"]
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(32)  # 폰트 크기 설정

        # 플래그 표시 (필요한 경우)
        if slide_data["flags"]:
            flag_text = ", ".join(slide_data["flags"])
            flag_textbox = slide.shapes.add_textbox(left, top - Inches(0.5), width, Inches(0.5))
            flag_textbox.text_frame.text = f"[{flag_text}]"

    prs.save("output.pptx")